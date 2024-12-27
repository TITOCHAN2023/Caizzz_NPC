import os
import pandas as pd
import pdfplumber
from docx import Document as DocxDocument
from bs4 import BeautifulSoup
from pptx import Presentation
import pytesseract
from PIL import Image

from langchain_core.documents.base import Document
from uuid import uuid4
from langchain_text_splitters import RecursiveCharacterTextSplitter
from env import SERVER

def extract_text_from_file(filepath):

    ext = os.path.splitext(filepath)[-1].lower()
    if not SERVER:
        return extract_text_from_file_cpu(filepath,ext)
    else:
        return extract_text_from_file_gpu(filepath,ext)
    

def extract_text_from_file_gpu(filepath,ext):


    from middleware.content_parser import pdfloader

    
    if ext == ".pdf":

        """HorYiny here u could help me (file could be png pdf ...)"""

    else:
        return extract_text_from_file_cpu(filepath,ext) #if can't use gpu, use cpu instead





def extract_text_from_file_cpu(filepath,ext):
    if ext == ".txt":
        with open(filepath, "r", encoding="utf-8") as f:
            return f.read()
        
    elif ext == ".pdf":
        with pdfplumber.open(filepath) as pdf:
            return "\n".join(page.extract_text() for page in pdf.pages)
        
    elif ext in [".docx",".doc"]:
        doc = DocxDocument(filepath)
        return "\n".join(para.text for para in doc.paragraphs)
    
    elif ext == ".html":
        with open(filepath, "r", encoding="utf-8") as f:
            soup = BeautifulSoup(f, "html.parser")
            return soup.get_text()
        
    elif ext in [".xlsx", ".xls"]:
        df = pd.read_excel(filepath)
        return df.to_string(index=False)
    
    elif ext == ".pptx":
        prs = Presentation(filepath)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        return "\n".join(text_runs)
    
    elif ext in [".png", ".jpg", ".jpeg"]:
        text = pytesseract.image_to_string(Image.open(filepath))
        
        return text
    else:
        logger.info(f"Unsupported file type: {ext}")
        return None



def load_and_split_documents(file_path):
    content = extract_text_from_file(file_path)
    text_splitter = RecursiveCharacterTextSplitter(
        chunk_size=500,
        chunk_overlap=50
    )
    texts = text_splitter.split_text(content)
    documents = [Document(page_content=t, metadata={"source": file_path}) for t in texts]

    return documents

